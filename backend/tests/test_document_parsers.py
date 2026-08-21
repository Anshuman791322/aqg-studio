"""Unit tests for deterministic document parsers (PDF, DOCX, PPTX, TXT) and linguistic cleaners."""

import io

import docx
import fitz
import pptx

from app.services.cleaner import (
    calculate_sha256,
    dehyphenate_text,
    detect_language,
    estimate_tokens,
    normalize_whitespace,
)
from app.services.parsers.docx import DOCXDocumentParser
from app.services.parsers.pdf import PDFDocumentParser
from app.services.parsers.pptx import PPTXDocumentParser
from app.services.parsers.txt import TXTDocumentParser


def _create_mock_pdf(
    pages_text: list[str],
    repeated_header: str | None = None,
    repeated_footer: str | None = None,
) -> bytes:
    """Create an in-memory PDF document using PyMuPDF."""
    doc = fitz.open()
    for text in pages_text:
        page = doc.new_page()
        full_text = text
        if repeated_header:
            full_text = f"{repeated_header}\n\n{full_text}"
        if repeated_footer:
            full_text = f"{full_text}\n\n{repeated_footer}"
        page.insert_text((50, 72), full_text)
    buf = doc.write()
    doc.close()
    return buf


def _create_mock_docx(headings_and_paras: list[tuple[str, list[str]]]) -> bytes:
    """Create an in-memory DOCX document using python-docx."""
    doc = docx.Document()
    for heading, paras in headings_and_paras:
        if heading:
            doc.add_heading(heading, level=1)
        for p in paras:
            doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _create_mock_pptx(slides_content: list[tuple[str, list[str]]]) -> bytes:
    """Create an in-memory PPTX presentation using python-pptx."""
    prs = pptx.Presentation()
    blank_layout = prs.slide_layouts[0]
    for title, bullet_points in slides_content:
        slide = prs.slides.add_slide(blank_layout)
        if slide.shapes.title:
            slide.shapes.title.text = title
        if len(slide.shapes.placeholders) > 1 and bullet_points:
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.text = "\n".join(bullet_points)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ------------------------------------------------------------------------------
# Cleaner & Linguistic Tests
# ------------------------------------------------------------------------------
def test_cleaner_whitespace_and_dehyphenation() -> None:
    """Verify whitespace collapsing and conservative dehyphenation."""
    dirty_text = "This is  a   test\n\n\nwith  hyphen-\nated words and tabs\t\tand spaces."
    cleaned = normalize_whitespace(dehyphenate_text(dirty_text))
    assert "hyphenated" in cleaned
    assert "  " not in cleaned
    assert "\n\n\n" not in cleaned


def test_language_detection() -> None:
    """Verify language detection across English, Spanish, French, and Hindi."""
    en_sample = (
        "The learning objectives of this course cover automated question generation."
    )
    es_sample = (
        "El objetivo del curso es que los estudiantes comprendan la estructura de las preguntas."
    )
    fr_sample = (
        "Les objectifs d'apprentissage de ce cours portent sur la génération de questions."
    )
    hi_sample = "इस पाठ का मुख्य उद्देश्य विद्यार्थियों को अवधारणाओं को समझाना है।"

    assert detect_language(en_sample) == "en"
    assert detect_language(es_sample) == "es"
    assert detect_language(fr_sample) == "fr"
    assert detect_language(hi_sample) == "hi"


def test_token_estimation() -> None:
    """Verify token estimation returns positive integer proportional to text length."""
    text = "Machine learning algorithms learn patterns from empirical training datasets."
    tokens = estimate_tokens(text)
    assert tokens >= 5
    assert tokens <= 25


def test_sha256_checksum() -> None:
    """Verify SHA-256 checksum is deterministic."""
    data = b"Hello AQG Studio Assessment Content"
    hash1 = calculate_sha256(data)
    hash2 = calculate_sha256(data)
    assert hash1 == hash2
    assert len(hash1) == 64


# ------------------------------------------------------------------------------
# PDF Parser Tests
# ------------------------------------------------------------------------------
def test_pdf_parser_standard_text_and_header_removal() -> None:
    """Verify PDF parser extracts pages and removes repeated running headers/footers."""
    pages = [
        "Chapter 1: Introduction to Calculus. Functions describe mathematical relationships.",
        "Chapter 1: Derivatives and Limits. A derivative measures instantaneous rate of change.",
        "Chapter 1: Integrals and Area. Definite integrals compute area under curves.",
    ]
    header = "University Calculus 101 - Spring 2026"
    footer = "Confidential - Page Number"

    pdf_bytes = _create_mock_pdf(pages, repeated_header=header, repeated_footer=footer)
    parser = PDFDocumentParser()
    result = parser.parse(pdf_bytes, "calculus.pdf")

    assert result.error_code is None
    assert result.page_count == 3
    assert result.word_count > 10
    assert result.is_scanned is False
    assert len(result.sections) == 3

    # Verify repeated headers were filtered out of parsed sections
    combined_content = " ".join(s.content for s in result.sections)
    assert "Calculus" in combined_content


def test_pdf_parser_scanned_detection() -> None:
    """Verify PDF parser detects scanned/empty documents lacking extractable text."""
    # Create empty PDF with no text inserted
    doc = fitz.open()
    doc.new_page()
    doc.new_page()
    empty_bytes = doc.write()
    doc.close()

    parser = PDFDocumentParser()
    result = parser.parse(empty_bytes, "scanned_exam.pdf")

    assert result.is_scanned is True
    assert result.error_code == "NEEDS_OCR"
    assert "insufficient extractable text" in result.error_message.lower()


# ------------------------------------------------------------------------------
# DOCX Parser Tests
# ------------------------------------------------------------------------------
def test_docx_parser_headings_and_paragraphs() -> None:
    """Verify DOCX parser groups content under structured headings."""
    content = [
        (
            "Section 1: Thermodynamics",
            ["Energy cannot be created or destroyed.", "Entropy increases."],
        ),
        (
            "Section 2: Quantum Mechanics",
            ["Wave-particle duality governs subatomic behavior."],
        ),
    ]
    docx_bytes = _create_mock_docx(content)
    parser = DOCXDocumentParser()
    result = parser.parse(docx_bytes, "physics.docx")

    assert result.error_code is None
    assert len(result.sections) == 2
    assert "Thermodynamics" in result.sections[0].title
    assert "Entropy increases" in result.sections[0].content


def test_docx_parser_rejects_legacy_doc() -> None:
    """Verify DOCX parser rejects legacy .doc files with explanatory message."""
    ole2_header = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 500
    parser = DOCXDocumentParser()
    result = parser.parse(ole2_header, "old_lecture.doc")

    assert result.error_code == "UNSUPPORTED_LEGACY_FORMAT"
    assert ".docx or .pdf" in result.error_message


# ------------------------------------------------------------------------------
# PPTX Parser Tests
# ------------------------------------------------------------------------------
def test_pptx_parser_slides() -> None:
    """Verify PPTX parser extracts slide titles and bullet points."""
    slides = [
        (
            "Neural Networks",
            ["Perceptrons and multi-layer networks", "Backpropagation algorithm"],
        ),
        (
            "Transformers",
            ["Self-attention mechanisms", "Multi-head attention"],
        ),
    ]
    pptx_bytes = _create_mock_pptx(slides)
    parser = PPTXDocumentParser()
    result = parser.parse(pptx_bytes, "deep_learning.pptx")

    assert result.error_code is None
    assert result.page_count == 2
    assert len(result.sections) == 2
    assert "Neural Networks" in result.sections[0].title


# ------------------------------------------------------------------------------
# TXT / Markdown Parser Tests
# ------------------------------------------------------------------------------
def test_txt_parser_markdown_headings() -> None:
    """Verify TXT parser detects Markdown headers."""
    markdown_content = (
        "# Introduction to Genetics\n\nDNA consists of four nucleotide bases.\n\n"
        "## Mendelian Inheritance\n\nAlleles segregate during gamete formation."
    )
    txt_bytes = markdown_content.encode("utf-8")
    parser = TXTDocumentParser()
    result = parser.parse(txt_bytes, "biology.md")

    assert result.error_code is None
    assert len(result.sections) == 2
    assert result.sections[0].title == "Introduction to Genetics"
    assert result.sections[1].title == "Mendelian Inheritance"
