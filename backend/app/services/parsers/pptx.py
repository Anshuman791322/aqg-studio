"""Deterministic PPTX presentation parser using python-pptx."""

import io
import zipfile

import pptx

from app.services.cleaner import (
    dehyphenate_text,
    detect_language,
    normalize_whitespace,
)
from app.services.parsers.base import BaseDocumentParser, ParsedDocument, ParsedSection


class PPTXDocumentParser(BaseDocumentParser):
    """Parser extracting text, slide titles, and speaker notes from PPTX presentations."""

    def parse(self, content_bytes: bytes, filename: str) -> ParsedDocument:
        # Validate PK zip signature
        if not content_bytes.startswith(b"PK"):
            return ParsedDocument(
                error_code="INVALID_FILE_SIGNATURE",
                error_message="The uploaded file does not contain a valid PPTX archive signature.",
            )

        # Protect against zip bombs / excessive entries / corrupted archives
        try:
            with zipfile.ZipFile(io.BytesIO(content_bytes)) as zf:
                entries = zf.infolist()
                if len(entries) > 1000:
                    return ParsedDocument(
                        error_code="ZIP_BOMB_DETECTED",
                        error_message="Presentation archive contains too many files.",
                    )

                total_uncompressed = sum(info.file_size for info in entries)
                if total_uncompressed > 200 * 1024 * 1024:
                    return ParsedDocument(
                        error_code="ZIP_BOMB_DETECTED",
                        error_message="Presentation exceeds uncompressed security size limits.",
                    )

                if total_uncompressed / max(len(content_bytes), 1) > 100:  # 100:1 ratio limit
                    return ParsedDocument(
                        error_code="ZIP_BOMB_DETECTED",
                        error_message="Presentation exceeds maximum safe compression ratio.",
                    )
        except Exception as e:
            return ParsedDocument(
                error_code="CORRUPTED_FILE",
                error_message=f"Failed to read PPTX archive: {str(e)}",
            )

        try:
            prs = pptx.Presentation(io.BytesIO(content_bytes))
        except Exception as e:
            return ParsedDocument(
                error_code="CORRUPTED_FILE",
                error_message=f"Failed to parse PPTX presentation: {str(e)}",
            )

        slide_count = len(prs.slides)
        if slide_count == 0:
            return ParsedDocument(
                page_count=0,
                error_code="EMPTY_DOCUMENT",
                error_message="The uploaded PPTX presentation contains no slides.",
            )

        sections: list[ParsedSection] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_title = f"Slide {slide_idx}"
            slide_texts: list[str] = []

            # Check for slide title shape
            if slide.shapes.title and slide.shapes.title.text:
                title_text = normalize_whitespace(slide.shapes.title.text)
                if title_text:
                    slide_title = f"Slide {slide_idx}: {title_text}"

            for shape in slide.shapes:
                # Avoid duplicating the title shape text
                if shape == slide.shapes.title:
                    continue

                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = normalize_whitespace(dehyphenate_text(paragraph.text))
                        if text:
                            slide_texts.append(text)

                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [
                            normalize_whitespace(cell.text)
                            for cell in row.cells
                            if normalize_whitespace(cell.text)
                        ]
                        if cells:
                            slide_texts.append(" | ".join(cells))

            # Check for speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = normalize_whitespace(slide.notes_slide.notes_text_frame.text)
                if notes_text:
                    slide_texts.append(f"Speaker Notes: {notes_text}")

            content = "\n\n".join(slide_texts)
            if content:
                sections.append(
                    ParsedSection(
                        title=slide_title,
                        content=content,
                        page_start=slide_idx,
                        page_end=slide_idx,
                        level=1,
                        metadata={"slide_number": slide_idx},
                    )
                )

        full_raw_text = "\n\n".join(f"{s.title}\n{s.content}" for s in sections)
        word_count = len(full_raw_text.split())
        language = detect_language(full_raw_text)

        return ParsedDocument(
            sections=sections,
            page_count=slide_count,
            word_count=word_count,
            language=language,
            raw_text=full_raw_text,
            is_scanned=False,
            metadata={"parser": "python-pptx", "slide_count": slide_count},
        )
